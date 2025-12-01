"""
内容解析器

从各种文件格式中提取文本内容，支持文档、图片、音视频等多种格式。
"""

import os
import re
from pathlib import Path
from typing import Optional, Dict, Any, List
import logging
from dataclasses import dataclass
from datetime import datetime
import asyncio

# 全局PaddleOCR实例，避免重复初始化
_paddle_ocr_instance = None
_paddle_ocr_lock = asyncio.Lock()

# 导入统一配置
try:
    from app.core.config import get_settings
    settings = get_settings()

    def get_parser_method(extension: str) -> str:
        return settings.mvp.get_parser_method(extension)

    def get_content_config(file_type: str) -> Dict[str, Any]:
        return settings.mvp.get_content_config(file_type)

    def is_mvp_mode() -> bool:
        return settings.mvp.is_mvp_mode()

    def get_format_display_name(extension: str) -> str:
        return settings.mvp.get_format_display_name(extension)
except ImportError:
    # 如果配置文件不存在，使用默认配置
    def get_parser_method(extension: str) -> str:
        return ''

    def get_content_config(file_type: str) -> Dict[str, Any]:
        return {}

    def is_mvp_mode() -> bool:
        return False

    def get_format_display_name(extension: str) -> str:
        return extension.upper()

# 文件处理库导入
try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False
    logging.warning("chardet未安装，编码检测功能不可用")

try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("PyPDF2未安装，PDF内容提取功能不可用")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx未安装，Word文档内容提取功能不可用")

try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logging.warning("openpyxl未安装，Excel文档内容提取功能不可用")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx未安装，PowerPoint文档内容提取功能不可用")

logger = logging.getLogger(__name__)


@dataclass
class ParsedContent:
    """解析后的内容数据类"""
    text: str
    title: Optional[str] = None
    language: Optional[str] = None
    encoding: Optional[str] = None
    confidence: float = 0.0
    metadata: Optional[Dict[str, Any]] = None


class ContentParser:
    """内容解析器

    支持从各种文件格式中提取文本内容：
    MVP阶段支持：
    - 文档类：PDF、Word、Excel、PowerPoint、TXT、Markdown
    - 音视频类：仅提取元数据（内容提取为未来功能）
    非MVP阶段支持：
    - 文本类：TXT、Markdown、代码文件
    - 预留接口：音视频转文字、图片OCR
    """

    def __init__(self, max_content_length: int = 1024 * 1024):  # 1MB
        """初始化内容解析器

        Args:
            max_content_length: 最大内容长度限制（字符数）
        """
        self.max_content_length = max_content_length
        self.mvp_mode = is_mvp_mode()

        if self.mvp_mode:
            # MVP模式：只支持PRD要求的格式
            self.supported_formats = {
                # Office文档解析 (现代格式 + 经典格式)
                '.pdf': self._parse_pdf,
                '.docx': self._parse_docx,
                '.xlsx': self._parse_excel,
                '.pptx': self._parse_pptx,
                '.doc': self._parse_doc,  # 经典Word格式
                '.xls': self._parse_excel,  # 经典Excel格式
                '.ppt': self._parse_ppt,   # 经典PowerPoint格式
                # 文本文档解析
                '.txt': self._parse_text,
                '.md': self._parse_markdown,
                # 音视频元数据解析
                '.mp3': self._parse_audio_metadata,
                '.wav': self._parse_audio_metadata,
                '.mp4': self._parse_video_metadata,
                '.avi': self._parse_video_metadata,
                # 图片内容解析
                '.png': self._parse_image_content,
                '.jpg': self._parse_image_content,
                '.jpeg': self._parse_image_content,
            }
            logger.info("使用MVP模式，支持PRD要求的核心文件格式")
        else:
            # 完整模式：支持所有格式
            self.supported_formats = {
                # Office文档解析 (现代格式 + 经典格式)
                '.pdf': self._parse_pdf,
                '.docx': self._parse_docx,
                '.xlsx': self._parse_excel,
                '.pptx': self._parse_pptx,
                '.doc': self._parse_doc,  # 经典Word格式
                '.xls': self._parse_excel,  # 经典Excel格式
                '.ppt': self._parse_ppt,   # 经典PowerPoint格式
                # 文本文档解析
                '.txt': self._parse_text,
                '.md': self._parse_markdown,
                '.rtf': self._parse_text,  # 简化处理
                # 代码文件解析
                '.py': self._parse_code,
                '.js': self._parse_code,
                '.ts': self._parse_code,
                '.html': self._parse_html,
                '.css': self._parse_code,
                '.java': self._parse_code,
                '.cpp': self._parse_code,
                '.c': self._parse_code,
                '.go': self._parse_code,
                '.rs': self._parse_code,
                '.php': self._parse_code,
                '.rb': self._parse_code,
                '.swift': self._parse_code,
                '.kt': self._parse_code,
                # 图片内容解析
                '.png': self._parse_image_content,
                '.jpg': self._parse_image_content,
                '.jpeg': self._parse_image_content,
            }
            logger.info("使用完整模式，支持所有文件格式")

        # 编码检测优先级
        self.encoding_priority = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1', 'ascii']

    async def parse_content(self, file_path: str) -> ParsedContent:
        """解析文件内容

        Args:
            file_path: 文件路径

        Returns:
            ParsedContent: 解析后的内容
        """
        try:
            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundError(f"文件不存在: {file_path}")

            # 根据文件扩展名选择解析方法
            extension = path.suffix.lower()
            parser_func = self.supported_formats.get(extension)

            if not parser_func:
                logger.warning(f"不支持的文件格式: {extension}")
                return ParsedContent(
                    text="",
                    confidence=0.0,
                    metadata={"error": f"不支持的文件格式: {extension}"}
                )

            # 执行解析 - 支持异步方法
            if extension in ['.mp3', '.wav', '.mp4', '.avi']:
                # 音视频文件需要异步处理
                if extension in ['.mp3', '.wav']:
                    parsed_content = await self._extract_audio_content(path)
                else:
                    parsed_content = await self._extract_video_content(path)
            elif extension in ['.png', '.jpg', '.jpeg']:
                # 图片文件需要异步处理
                parsed_content = await self._extract_image_content(path)
            else:
                # 其他文件类型使用同步处理
                parsed_content = parser_func(path)

            # 内容长度限制
            if len(parsed_content.text) > self.max_content_length:
                parsed_content.text = parsed_content.text[:self.max_content_length]
                parsed_content.text += "\n... [内容被截断]"

            # 语言检测（简化版）
            if not parsed_content.language:
                parsed_content.language = self._detect_language(parsed_content.text)

            return parsed_content

        except Exception as e:
            logger.error(f"解析文件内容失败 {file_path}: {e}")
            return ParsedContent(text="", confidence=0.0, metadata={"error": str(e)})

    def _parse_pdf(self, path: Path) -> ParsedContent:
        """解析PDF文件内容"""
        if not PDF_AVAILABLE:
            return ParsedContent(text="", error="PDF解析库不可用")

        try:
            text_parts = []
            title = None

            with open(path, 'rb') as file:
                reader = PdfReader(file)

                # 尝试提取标题
                if reader.metadata and reader.metadata.get('/Title'):
                    title = reader.metadata.get('/Title')

                # 提取每一页的文本
                for page_num, page in enumerate(reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            # 清理PDF文本中的乱码和格式问题
                            cleaned_text = self._clean_pdf_text(page_text)
                            if cleaned_text.strip():
                                text_parts.append(f"[页面 {page_num + 1}]\n{cleaned_text}")
                    except Exception as e:
                        logger.warning(f"提取PDF第{page_num + 1}页内容失败: {e}")

            text = "\n\n".join(text_parts)

            return ParsedContent(
                text=text,
                title=title,
                language=self._detect_language(text),
                confidence=0.8 if text else 0.0
            )

        except Exception as e:
            logger.error(f"解析PDF文件失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _parse_docx(self, path: Path) -> ParsedContent:
        """解析Word文档内容"""
        if not DOCX_AVAILABLE:
            return ParsedContent(text="", error="Word文档解析库不可用")

        try:
            doc = Document(str(path))

            # 提取标题
            title = None
            if doc.core_properties.title:
                title = doc.core_properties.title
            elif doc.paragraphs and doc.paragraphs[0].text.strip():
                # 使用第一段作为标题
                title = doc.paragraphs[0].text.strip()

            # 提取段落文本
            paragraph_texts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraph_texts.append(para.text)

            # 提取表格内容
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_texts.append(" | ".join(row_text))

            # 组合文本
            text_parts = []
            if paragraph_texts:
                text_parts.extend(paragraph_texts)
            if table_texts:
                text_parts.append("\n[表格内容]\n" + "\n".join(table_texts))

            text = "\n\n".join(text_parts)

            return ParsedContent(
                text=text,
                title=title,
                language=self._detect_language(text),
                confidence=0.9 if text else 0.0
            )

        except Exception as e:
            logger.error(f"解析Word文档失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _parse_excel(self, path: Path) -> ParsedContent:
        """解析Excel文件内容"""
        if not EXCEL_AVAILABLE:
            return ParsedContent(text="", error="Excel解析库不可用")

        try:
            # 根据文件大小决定读取模式
            file_size = path.stat().st_size
            use_read_only = file_size > 50 * 1024 * 1024  # 50MB以上使用read_only模式

            if use_read_only:
                workbook = load_workbook(str(path), read_only=True)
            else:
                workbook = load_workbook(str(path), read_only=False, data_only=True)

            text_parts = []
            total_data_rows = 0

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []

                # 检查工作表大小
                max_row = sheet.max_row if hasattr(sheet, 'max_row') else 1

                # 限制处理的最大行数，防止内存溢出
                max_process_rows = min(1000, max_row)  # 最多处理1000行

                for row_idx, row in enumerate(sheet.iter_rows(values_only=True, max_row=max_process_rows)):
                    # 过滤空行
                    if any(cell is not None and str(cell).strip() for cell in row):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                cell_str = str(cell).strip()
                                if cell_str:
                                    row_text.append(cell_str)
                        if row_text:
                            sheet_data.append(" | ".join(row_text))
                            total_data_rows += 1

                    # 如果处理了足够的行，提前结束
                    if row_idx >= max_process_rows - 1:
                        break

                # 添加工作表信息
                if sheet_data:
                    header = f"[工作表: {sheet_name} (显示前{len(sheet_data)}行数据)]"
                    text_parts.append(header + "\n" + "\n".join(sheet_data))

                    # 如果有更多数据未显示，添加说明
                    if max_row > max_process_rows:
                        remaining_rows = max_row - max_process_rows
                        text_parts.append(f"\n... (还有 {remaining_rows:,} 行数据未显示)")
                else:
                    # 检查是否有表头等基础信息
                    try:
                        first_row = next(sheet.iter_rows(values_only=True), None)
                        if first_row and any(cell is not None and str(cell).strip() for cell in first_row):
                            header_text = " | ".join(str(cell).strip() for cell in first_row if cell and str(cell).strip())
                            text_parts.append(f"[工作表: {sheet_name} - 表头]\n" + header_text)
                            total_data_rows += 1
                    except:
                        pass

            text = "\n\n".join(text_parts)

            return ParsedContent(
                text=text if text else f"[Excel文档: {path.name} - 暂无可提取的文本内容]",
                title=f"Excel文档 - {path.name}",
                language=self._detect_language(text) if text else "zh",
                confidence=0.8 if total_data_rows > 0 else 0.3
            )

        except Exception as e:
            logger.error(f"解析Excel文件失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _parse_pptx(self, path: Path) -> ParsedContent:
        """解析PowerPoint文档内容"""
        if not PPTX_AVAILABLE:
            return ParsedContent(text="", error="PowerPoint解析库不可用")

        try:
            prs = Presentation(str(path))
            slide_texts = []

            for slide_num, slide in enumerate(prs.slides):
                slide_content = []

                # 提取幻灯片标题和内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 清理PPTX文本中的乱码字符（垂直制表符等）
                        cleaned_text = self._clean_pptx_text(shape.text)
                        if cleaned_text.strip():
                            slide_content.append(cleaned_text.strip())

                if slide_content:
                    slide_texts.append(f"[幻灯片 {slide_num + 1}]\n" + "\n".join(slide_content))

            text = "\n\n".join(slide_texts)

            # 提取标题
            title = None
            if prs.core_properties.title:
                title = prs.core_properties.title
            elif slide_texts:
                # 使用第一张幻灯片的内容作为标题
                first_slide = slide_texts[0].split('\n', 1)[1].split('\n')[0] if '\n' in slide_texts[0] else ""
                title = first_slide.strip()

            return ParsedContent(
                text=text,
                title=title,
                language=self._detect_language(text),
                confidence=0.8 if text else 0.0
            )

        except Exception as e:
            logger.error(f"解析PowerPoint文档失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _parse_text(self, path: Path) -> ParsedContent:
        """解析纯文本文件"""
        try:
            content, encoding = self._read_text_file(path)

            if not content:
                return ParsedContent(text="", title=path.name, confidence=0.0)

            # 尝试提取标题（第一行）
            lines = content.split('\n')
            title = lines[0].strip() if lines and lines[0].strip() else path.name

            return ParsedContent(
                text=content,
                title=title,
                language=self._detect_language(content),
                encoding=encoding,
                confidence=0.9 if content else 0.0
            )

        except Exception as e:
            logger.error(f"解析文本文件失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _parse_markdown(self, path: Path) -> ParsedContent:
        """解析Markdown文件"""
        try:
            content, encoding = self._read_text_file(path)

            if not content:
                return ParsedContent(text="", title=path.name, confidence=0.0)

            # 提取Markdown标题
            title = self._extract_markdown_title(content)

            # 清理Markdown语法
            clean_text = self._clean_markdown(content)

            return ParsedContent(
                text=clean_text,
                title=title,
                language=self._detect_language(content),
                encoding=encoding,
                confidence=0.9 if content else 0.0
            )

        except Exception as e:
            logger.error(f"解析Markdown文件失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _parse_code(self, path: Path) -> ParsedContent:
        """解析代码文件"""
        try:
            content, encoding = self._read_text_file(path)

            if not content:
                return ParsedContent(text="", title=path.name, confidence=0.0)

            # 提取注释和文档字符串
            code_content = self._extract_code_comments(path.suffix, content)

            return ParsedContent(
                text=code_content,
                title=path.name,
                language=self._detect_language(content),
                encoding=encoding,
                confidence=0.7 if code_content else 0.0
            )

        except Exception as e:
            logger.error(f"解析代码文件失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _parse_html(self, path: Path) -> ParsedContent:
        """解析HTML文件"""
        try:
            content, encoding = self._read_text_file(path)

            if not content:
                return ParsedContent(text="", title=path.name, confidence=0.0)

            # 简单的HTML标签清理
            clean_text = re.sub(r'<[^>]+>', ' ', content)
            clean_text = re.sub(r'\s+', ' ', clean_text).strip()

            # 提取title标签
            title_match = re.search(r'<title[^>]*>(.*?)</title>', content, re.IGNORECASE | re.DOTALL)
            title = title_match.group(1).strip() if title_match else path.name

            return ParsedContent(
                text=clean_text,
                title=title,
                language=self._detect_language(clean_text),
                encoding=encoding,
                confidence=0.6 if clean_text else 0.0
            )

        except Exception as e:
            logger.error(f"解析HTML文件失败 {path}: {e}")
            return ParsedContent(text="", error=str(e))

    def _read_text_file(self, path: Path) -> tuple[str, Optional[str]]:
        """读取文本文件，自动检测编码"""
        encodings_to_try = self.encoding_priority.copy()

        # 如果有chardet，先检测编码
        if CHARDET_AVAILABLE:
            try:
                with open(path, 'rb') as f:
                    raw_data = f.read(10240)  # 读取前10KB进行检测
                result = chardet.detect(raw_data)
                if result and result['confidence'] > 0.7:
                    detected_encoding = result['encoding']
                    if detected_encoding and detected_encoding not in encodings_to_try:
                        encodings_to_try.insert(0, detected_encoding)
            except Exception as e:
                logger.warning(f"编码检测失败 {path}: {e}")

        # 尝试不同编码读取文件
        for encoding in encodings_to_try:
            try:
                with open(path, 'r', encoding=encoding) as f:
                    content = f.read()
                return content, encoding
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"读取文件失败 {path} (编码: {encoding}): {e}")
                continue

        # 如果所有编码都失败，使用默认编码并忽略错误
        try:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return content, 'utf-8'
        except Exception as e:
            logger.error(f"读取文件完全失败 {path}: {e}")
            return "", None

    def _extract_markdown_title(self, content: str) -> str:
        """从Markdown内容中提取标题"""
        lines = content.split('\n')

        # 查找一级标题
        for line in lines:
            line = line.strip()
            if line.startswith('# '):
                return line[2:].strip()

        # 如果没有一级标题，使用第一行非空文本
        for line in lines:
            if line.strip():
                return line.strip()

        return ""

    def _clean_markdown(self, content: str) -> str:
        """清理Markdown语法，保留纯文本"""
        # 移除代码块
        content = re.sub(r'```.*?```', '', content, flags=re.DOTALL)
        # 移除行内代码
        content = re.sub(r'`([^`]+)`', r'\1', content)
        # 移除标题标记
        content = re.sub(r'^#+\s*', '', content, flags=re.MULTILINE)
        # 移除粗体和斜体标记
        content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
        content = re.sub(r'\*([^*]+)\*', r'\1', content)
        # 移除链接
        content = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', content)
        # 移除图片
        content = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'\1', content)
        # 移除列表标记
        content = re.sub(r'^\s*[-*+]\s*', '', content, flags=re.MULTILINE)
        content = re.sub(r'^\s*\d+\.\s*', '', content, flags=re.MULTILINE)
        # 清理多余空白
        content = re.sub(r'\n\s*\n', '\n\n', content)

        return content.strip()

    def _extract_code_comments(self, extension: str, content: str) -> str:
        """从代码中提取注释和文档字符串"""
        lines = content.split('\n')
        comments = []

        # 不同语言的注释模式
        comment_patterns = {
            '.py': [
                r'^\s*"""[^"]*"""',  # 多行文档字符串
                r'^\s*\'\'\'[^\']*\'\'\'',
                r'^\s*#.*$',  # 单行注释
            ],
            '.js': [
                r'^\s*/\*.*?\*/',  # 多行注释
                r'^\s*//.*$',  # 单行注释
            ],
            '.ts': [
                r'^\s*/\*.*?\*/',
                r'^\s*//.*$',
            ],
            '.java': [
                r'^\s*/\*.*?\*/',
                r'^\s*//.*$',
            ],
            '.cpp': [
                r'^\s*/\*.*?\*/',
                r'^\s*//.*$',
            ],
            '.c': [
                r'^\s*/\*.*?\*/',
                r'^\s*//.*$',
            ],
            '.go': [
                r'^\s*/\*.*?\*/',
                r'^\s*//.*$',
            ],
            '.rs': [
                r'^\s*/\*.*?\*/',
                r'^\s*//.*$',
            ],
        }

        patterns = comment_patterns.get(extension, [r'^\s*//.*$', r'^\s*#.*$'])

        for line in lines:
            for pattern in patterns:
                if re.match(pattern, line, re.MULTILINE):
                    comment = re.sub(r'^\s*(/\*|\*/|\*|//|#)', '', line).strip()
                    if comment:
                        comments.append(comment)
                    break

        return '\n'.join(comments)

    def _clean_pdf_text(self, text: str) -> str:
        """清理PDF文本中的乱码和格式问题"""
        if not text:
            return ""

        import unicodedata

        # 移除控制字符和不可见字符
        cleaned_text = ''.join(char for char in text if unicodedata.category(char)[0] != 'C')

        # 识别和移除重复的乱码模式
        # 使用正则表达式匹配连续重复的相同乱码文本
        import re

        # 匹配连续重复的非正常文本模式（4次或以上重复）
        def remove_repeated_garbage(text):
            # 匹配4次或以上重复的相同文本片段
            pattern = r'((.{2,20})\2{3,})'
            return re.sub(pattern, '', text)

        cleaned_text = remove_repeated_garbage(cleaned_text)

        # 替换常见的PDF乱码字符
        replacements = {
            '�': '',  # 替换字符
            '□': '',  # 方框乱码
            '■': '',  # 实心方框
            '▪': '',  # 小方块
            '▫': '',  # 空心方块
            '▬': '',  # 长方形
            '▭': '',  # 空心长方形
            '©': '(c)',  # 版权符号
            '®': '(r)',  # 注册商标
            '™': '(tm)',  # 商标
        }

        for old, new in replacements.items():
            cleaned_text = cleaned_text.replace(old, new)

        # 移除包含大量非中文、非英文、非数字字符的行（可能是乱码行）
        lines = cleaned_text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            if line:
                # 计算有效字符比例（中文、英文、数字、标点符号）
                valid_chars = len(re.findall(r'[\u4e00-\u9fff\w\s.,;:!?()[\]{}"\'-]', line))
                total_chars = len(line)

                # 如果有效字符比例低于60%，可能是乱码行，跳过
                if total_chars == 0 or valid_chars / total_chars >= 0.6:
                    cleaned_lines.append(line)

        # 清理多余的空白字符
        cleaned_text = ' '.join(cleaned_lines)
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

        # 移除孤立的数字和字母（通常是PDF页码或页眉页脚）
        lines = cleaned_text.split('\n')
        final_lines = []
        for line in lines:
            line = line.strip()
            # 跳过纯数字行或非常短的行（可能是页码）
            if line and not (line.isdigit() or (len(line) <= 3 and line.isalnum())):
                final_lines.append(line)

        # 重新组合文本
        final_text = ' '.join(final_lines)

        # 最终清理多余空格和乱码模式
        final_text = re.sub(r'\s+', ' ', final_text).strip()

        # 移除开头和结尾的重复乱码文本
        final_text = re.sub(r'^[^\u4e00-\u9fff\w]{10,}', '', final_text)
        final_text = re.sub(r'[^\u4e00-\u9fff\w]{10,}$', '', final_text)

        return final_text

    def _clean_pptx_text(self, text: str) -> str:
        """清理PPTX文本中的乱码和格式问题"""
        if not text:
            return ""

        import re

        # 移除垂直制表符和其他常见控制字符
        text = text.replace('\x0B', '')  # 垂直制表符
        text = text.replace('\x0C', '')  # 换页符
        text = text.replace('\x08', '')  # 退格符
        text = text.replace('\x07', '')  # 响铃符
        text = text.replace('\x1B', '')  # ESC转义字符

        # 移除连续的特殊控制字符组合
        text = re.sub(r'[\x00-\x08\x0B-\x1F\x7F]', '', text)

        # 替换常见的PPTX乱码模式
        # 去除连续重复的特殊字符
        text = re.sub(r'[^\w\s\u4e00-\u9fff.,;:!?()[\]{}"\'\-]{4,}', '', text)

        # 清理多余的空白字符
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()

        return text

    def _detect_language(self, text: str) -> str:
        """简单的语言检测"""
        if not text or len(text) < 10:
            return "unknown"

        # 简单的中文检测
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
        english_chars = len(re.findall(r'[a-zA-Z]', text))

        if chinese_chars > english_chars:
            return "zh"
        elif english_chars > 0:
            return "en"
        else:
            return "unknown"

    async def _extract_audio_content(self, path: Path) -> ParsedContent:
        """从音频文件中提取语音内容"""
        try:
            import librosa
            import soundfile as sf
            import tempfile
            import os
        except ImportError as e:
            logger.warning(f"音频处理库不可用: {e}")
            return self._parse_audio_metadata_fallback(path)

        try:
            # 获取音频基本信息
            extension = path.suffix.lower()
            file_size = path.stat().st_size

            # 使用librosa获取音频时长
            duration = librosa.get_duration(path=str(path))

            # 处理15分钟时长限制
            max_duration = 15 * 60  # 15分钟 = 900秒
            if duration > max_duration:
                logger.info(f"音频文件时长超过限制: {duration:.1f}秒 > {max_duration}秒，将提取前{max_duration}秒内容")

                # 使用临时文件截取前15分钟
                try:
                    import tempfile
                    import subprocess

                    with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as temp_audio_file:
                        temp_audio_path = temp_audio_file.name

                    # 使用ffmpeg截取前15分钟音频
                    ffmpeg_cmd = [
                        'ffmpeg', '-i', str(path),
                        '-t', str(max_duration),  # 截取前15分钟
                        '-c', 'copy',  # 复制编解码器
                        '-y',  # 覆盖输出文件
                        temp_audio_path
                    ]

                    result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True, timeout=30)

                    if result.returncode != 0:
                        logger.warning(f"ffmpeg截取音频失败: {result.stderr}")
                        # 如果截取失败，继续使用原文件处理
                        audio_path_for_transcription = str(path)
                        truncated = False
                    else:
                        logger.info(f"成功截取音频前15分钟: {temp_audio_path}")
                        audio_path_for_transcription = temp_audio_path
                        truncated = True

                except Exception as e:
                    logger.warning(f"音频截取失败，使用原文件: {str(e)}")
                    audio_path_for_transcription = str(path)
                    truncated = False
            else:
                audio_path_for_transcription = str(path)
                truncated = False

            logger.info(f"开始处理音频文件: {path.name}, 时长: {duration:.1f}秒")

            # 获取AI模型服务进行语音转文字
            try:
                from app.services.ai_model_manager import ai_model_service

                # 调用Whisper模型进行语音识别
                transcription_result = await ai_model_service.speech_to_text(
                    audio_path_for_transcription,
                    language="zh"
                )

                if transcription_result and transcription_result.get("success", False):
                    transcribed_text = transcription_result.get("text", "").strip()
                    confidence = transcription_result.get("confidence", 0.0)

                    # 尝试提取元数据
                    metadata = await self._extract_audio_metadata_with_mutagen(path)
                    metadata.update({
                        "transcribed": True,
                        "transcription_confidence": confidence,
                        "original_duration": duration,
                        "file_size": file_size,
                        "truncated": truncated,
                        "processed_duration": min(duration, max_duration) if truncated else duration
                    })

                    # 提取标题（从元数据或文件名）
                    title = metadata.get('TIT2') or metadata.get('title') or path.stem

                    logger.info(f"音频转录完成: {path.name}, 文本长度: {len(transcribed_text)}字符")

                    return ParsedContent(
                        text=transcribed_text,
                        title=title,
                        language="zh",
                        confidence=confidence,
                        metadata=metadata
                    )
                else:
                    error_msg = transcription_result.get("error", "语音识别失败")
                    logger.warning(f"语音识别失败: {path.name}, 错误: {error_msg}")

                    # 降级为元数据提取
                    metadata = await self._extract_audio_metadata_with_mutagen(path)
                    metadata.update({
                        "transcribed": False,
                        "transcription_error": error_msg,
                        "original_duration": duration,
                        "file_size": file_size,
                        "truncated": truncated,
                        "processed_duration": min(duration, max_duration) if truncated else duration
                    })

                    return ParsedContent(
                        text=f"[语音识别失败: {error_msg}] - 仅提取元数据",
                        title=metadata.get('TIT2') or metadata.get('title') or path.stem,
                        language="metadata",
                        confidence=0.3,
                        metadata=metadata
                    )

            except Exception as e:
                logger.error(f"调用AI模型服务失败: {str(e)}")
                return self._parse_audio_metadata_fallback(path, duration)

        except Exception as e:
            logger.error(f"提取音频内容失败 {path}: {e}")
            return self._parse_audio_metadata_fallback(path)

        finally:
            # 清理临时截取文件
            if truncated and 'temp_audio_path' in locals():
                try:
                    if os.path.exists(temp_audio_path):
                        os.unlink(temp_audio_path)
                        logger.info(f"清理临时音频文件: {temp_audio_path}")
                except Exception as e:
                    logger.warning(f"清理临时文件失败: {str(e)}")

    async def _extract_audio_metadata_with_mutagen(self, path: Path) -> Dict[str, Any]:
        """使用mutagen提取音频元数据"""
        metadata = {
            "format": "audio",
            "file_extension": path.suffix.lower(),
            "file_size": path.stat().st_size
        }

        try:
            import mutagen
            from mutagen.mp3 import MP3
            from mutagen.wave import WAVE

            extension = path.suffix.lower()

            if extension == '.mp3':
                audio = MP3(str(path))
                if audio.info:
                    metadata.update({
                        "bitrate": audio.info.bitrate,
                        "sample_rate": audio.info.sample_rate,
                        "channels": audio.info.channels,
                        "layer": audio.info.layer,
                        "version": audio.info.version
                    })

                # 提取标签信息
                if audio.tags:
                    for key, value in audio.tags.items():
                        if isinstance(value, list) and len(value) == 1:
                            metadata[key] = str(value[0])
                        else:
                            metadata[key] = str(value)

            elif extension == '.wav':
                audio = WAVE(str(path))
                if audio.info:
                    metadata.update({
                        "sample_rate": audio.info.sample_rate,
                        "channels": audio.info.channels,
                        "bits_per_sample": getattr(audio.info, 'bits_per_sample', 0)
                    })

        except ImportError:
            logger.warning("mutagen库不可用，跳过详细元数据提取")
        except Exception as e:
            logger.warning(f"提取音频元数据失败: {e}")

        return metadata

    def _parse_audio_metadata(self, path: Path) -> ParsedContent:
        """解析音频文件元数据（同步方法，用于索引服务检查）"""
        try:
            format_name = get_format_display_name(path.suffix)
            metadata = {
                "format": "audio",
                "file_extension": path.suffix.lower(),
                "file_size": path.stat().st_size,
                "transcribed": False,
                "metadata_only": True
            }

            return ParsedContent(
                text=f"[{format_name} 音频文件] - 元数据已记录",
                title=path.stem,
                language="metadata",
                confidence=0.5,
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"音频元数据解析失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "audio", "error": str(e)}
            )

    def _parse_video_metadata(self, path: Path) -> ParsedContent:
        """解析视频文件元数据（同步方法，用于索引服务检查）"""
        try:
            format_name = get_format_display_name(path.suffix)
            metadata = {
                "format": "video",
                "file_extension": path.suffix.lower(),
                "file_size": path.stat().st_size,
                "transcribed": False,
                "metadata_only": True
            }

            return ParsedContent(
                text=f"[{format_name} 视频文件] - 元数据已记录",
                title=path.stem,
                language="metadata",
                confidence=0.5,
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"视频元数据解析失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "video", "error": str(e)}
            )

    def _parse_image_content(self, path: Path) -> ParsedContent:
        """解析图片文件内容（同步方法，用于索引服务检查）"""
        try:
            format_name = get_format_display_name(path.suffix)
            metadata = {
                "format": "image",
                "file_extension": path.suffix.lower(),
                "file_size": path.stat().st_size,
                "content_analyzed": False,
                "metadata_only": True
            }

            return ParsedContent(
                text=f"[{format_name} 图片文件] - 元数据已记录",
                title=path.stem,
                language="metadata",
                confidence=0.5,
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"图片内容解析失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "image", "error": str(e)}
            )

    def _parse_audio_metadata_fallback(self, path: Path, duration: float = None) -> ParsedContent:
        """音频解析降级方案：仅提取元数据"""
        try:
            format_name = get_format_display_name(path.suffix)
            metadata = {
                "format": "audio",
                "file_extension": path.suffix.lower(),
                "file_size": path.stat().st_size,
                "transcribed": False,
                "fallback": True
            }

            if duration is not None:
                metadata["duration"] = duration

            return ParsedContent(
                text=f"[{format_name} 元数据已提取] - 内容提取功能暂不可用",
                title=path.stem,
                language="metadata",
                confidence=0.6,
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"音频元数据提取失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "audio", "error": str(e)}
            )

    async def _extract_video_content(self, path: Path) -> ParsedContent:
        """从视频文件中提取音频内容"""
        try:
            import cv2
            import librosa
            import soundfile as sf
            import tempfile
            import os
            import subprocess
        except ImportError as e:
            logger.warning(f"视频处理库不可用: {e}")
            return self._parse_video_metadata_fallback(path)

        try:
            # 获取视频基本信息
            extension = path.suffix.lower()
            file_size = path.stat().st_size

            # 获取视频时长
            cap = cv2.VideoCapture(str(path))
            if not cap.isOpened():
                raise ValueError("无法打开视频文件")

            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))

            if fps > 0:
                duration = frame_count / fps
            else:
                duration = 0

            cap.release()

            # 处理15分钟时长限制
            max_duration = 15 * 60  # 15分钟 = 900秒
            if duration > max_duration:
                logger.info(f"视频文件时长超过限制: {duration:.1f}秒 > {max_duration}秒，将提取前{max_duration}秒内容")
                video_truncated = True
            else:
                video_truncated = False

            logger.info(f"开始处理视频文件: {path.name}, 时长: {duration:.1f}秒")

            # 提取音频轨道
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_audio_file:
                temp_audio_path = temp_audio_file.name

            try:
                # 动态构建ffmpeg命令
                ffmpeg_cmd = [
                    'ffmpeg', '-i', str(path),
                    '-vn',  # 不要视频
                    '-acodec', 'pcm_s16le',  # 音频编码
                    '-ar', '16000',  # 采样率16kHz
                    '-ac', '1',  # 单声道
                ]

                # 如果需要截取前15分钟，添加-t参数
                if video_truncated:
                    ffmpeg_cmd.extend(['-t', str(max_duration)])

                ffmpeg_cmd.extend(['-y', temp_audio_path])  # 覆盖输出文件

                result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True, timeout=min(duration, max_duration) + 60)

                if result.returncode != 0:
                    logger.warning(f"ffmpeg提取音频失败: {result.stderr}")
                    return self._parse_video_metadata_fallback(path, duration, width, height, fps)

                logger.info(f"音频轨道提取完成: {temp_audio_path} (截取: {video_truncated})")

                # 使用AI模型服务进行语音转文字
                try:
                    from app.services.ai_model_manager import ai_model_service

                    # 调用Whisper模型进行语音识别
                    transcription_result = await ai_model_service.speech_to_text(
                        temp_audio_path,
                        language="zh"
                    )

                    if transcription_result and transcription_result.get("success", False):
                        transcribed_text = transcription_result.get("text", "").strip()
                        confidence = transcription_result.get("confidence", 0.0)

                        logger.info(f"视频音频转录完成: {path.name}, 文本长度: {len(transcribed_text)}字符")

                        # 构建元数据
                        metadata = {
                            "format": "video",
                            "file_extension": extension,
                            "original_duration": duration,
                            "file_size": file_size,
                            "resolution": f"{width}x{height}",
                            "fps": fps,
                            "transcribed": True,
                            "transcription_confidence": confidence,
                            "audio_extracted": True,
                            "truncated": video_truncated,
                            "processed_duration": min(duration, max_duration) if video_truncated else duration
                        }

                        return ParsedContent(
                            text=transcribed_text,
                            title=path.stem,
                            language="zh",
                            confidence=confidence,
                            metadata=metadata
                        )
                    else:
                        error_msg = transcription_result.get("error", "语音识别失败")
                        logger.warning(f"视频语音识别失败: {path.name}, 错误: {error_msg}")

                        metadata = {
                            "format": "video",
                            "file_extension": extension,
                            "original_duration": duration,
                            "file_size": file_size,
                            "resolution": f"{width}x{height}",
                            "fps": fps,
                            "audio_extracted": True,
                            "transcribed": False,
                            "transcription_error": error_msg,
                            "truncated": video_truncated,
                            "processed_duration": min(duration, max_duration) if video_truncated else duration
                        }

                        return ParsedContent(
                            text=f"[语音识别失败: {error_msg}] - 音频轨道已提取",
                            title=path.stem,
                            language="metadata",
                            confidence=0.3,
                            metadata=metadata
                        )

                except Exception as e:
                    logger.error(f"调用AI模型服务失败: {str(e)}")
                    return self._parse_video_metadata_fallback(path, duration, width, height, fps)

            finally:
                # 清理临时音频文件
                if 'temp_audio_path' in locals() and os.path.exists(temp_audio_path):
                    try:
                        os.unlink(temp_audio_path)
                        logger.info(f"清理临时视频音频文件: {temp_audio_path}")
                    except Exception as e:
                        logger.warning(f"清理临时文件失败: {str(e)}")

        except subprocess.TimeoutExpired:
            logger.error(f"视频处理超时: {path}")
            return self._parse_video_metadata_fallback(path)
        except Exception as e:
            logger.error(f"提取视频内容失败 {path}: {e}")
            return self._parse_video_metadata_fallback(path)

    def _parse_video_metadata_fallback(self, path: Path, duration: float = None, width: int = None, height: int = None, fps: float = None) -> ParsedContent:
        """视频解析降级方案：仅提取元数据"""
        try:
            format_name = get_format_display_name(path.suffix)
            metadata = {
                "format": "video",
                "file_extension": path.suffix.lower(),
                "file_size": path.stat().st_size,
                "transcribed": False,
                "fallback": True
            }

            if duration is not None:
                metadata["duration"] = duration
            if width is not None and height is not None:
                metadata["resolution"] = f"{width}x{height}"
            if fps is not None:
                metadata["fps"] = fps

            return ParsedContent(
                text=f"[{format_name} 元数据已提取] - 内容提取功能暂不可用",
                title=path.stem,
                language="metadata",
                confidence=0.6,
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"视频元数据提取失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "video", "error": str(e)}
            )

    async def _extract_image_content(self, path: Path) -> ParsedContent:
        """从图片文件中提取内容"""
        try:
            from PIL import Image
            import numpy as np
        except ImportError as e:
            logger.warning(f"图片处理库不可用: {e}")
            return self._parse_image_metadata_fallback(path)

        try:
            # 获取图片基本信息
            extension = path.suffix.lower()
            file_size = path.stat().st_size

            logger.info(f"开始处理图片文件: {path.name}, 大小: {file_size}字节")

            # 使用OCR识别图片中的文字内容
            try:
                # 提取图片文字内容
                ocr_text = await self._extract_text_from_image(path)

                # 构建图片描述信息
                description_parts = []

                # 基本信息
                description_parts.append(f"{extension.upper()}格式图片")

                # 添加OCR识别的文字内容
                if ocr_text.strip():
                    description_parts.append(f"文字内容：{ocr_text}")
                    # OCR成功，置信度较高
                    confidence = 0.9
                    ocr_success = True
                else:
                    description_parts.append("未检测到文字内容")
                    confidence = 0.6  # 只有基本信息，置信度较低
                    ocr_success = False

                # 组合描述文本
                image_description = " | ".join(description_parts)

                # 提取图片元数据
                with Image.open(path) as img:
                    width, height = img.size
                    metadata = {
                        "format": "image",
                        "file_extension": extension,
                        "file_size": file_size,
                        "width": width,
                        "height": height,
                        "mode": img.mode,
                        "ocr_extracted": ocr_success,
                        "ocr_text_length": len(ocr_text.strip()),
                        "processed_at": datetime.now().isoformat()
                    }

                logger.info(f"图片OCR处理完成: {path.name}, 文字长度: {len(ocr_text)}字符")

                return ParsedContent(
                    text=image_description,
                    title=path.stem,
                    language="zh",
                    confidence=confidence,
                    metadata=metadata
                )

            except Exception as e:
                logger.error(f"OCR处理图片失败 {path}: {str(e)}")
                # OCR失败时降级为元数据提取
                return self._parse_image_metadata_fallback(path, f"OCR处理失败: {str(e)}")

        except Exception as e:
            logger.error(f"提取图片内容失败 {path}: {e}")
            return self._parse_image_metadata_fallback(path)

    def _parse_image_metadata_fallback(self, path: Path, error_msg: str = None) -> ParsedContent:
        """图片解析降级方案：仅提取元数据"""
        try:
            from PIL import Image

            format_name = get_format_display_name(path.suffix)
            metadata = {
                "format": "image",
                "file_extension": path.suffix.lower(),
                "file_size": path.stat().st_size,
                "image_understood": False,
                "fallback": True
            }

            # 尝试获取图片基本信息
            try:
                with Image.open(path) as img:
                    metadata.update({
                        "width": img.width,
                        "height": img.height,
                        "mode": img.mode
                    })
            except Exception as e:
                logger.warning(f"无法读取图片基本信息: {e}")

            # 构建错误消息
            if error_msg:
                text = f"[图像理解失败: {error_msg}] - 仅提取元数据"
            else:
                text = f"[{format_name} 元数据已提取] - 内容提取功能暂不可用"

            return ParsedContent(
                text=text,
                title=path.stem,
                language="metadata",
                confidence=0.3,
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"图片元数据提取失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "image", "error": str(e)}
            )

    def _parse_doc(self, path: Path) -> ParsedContent:
        """解析经典Word文档 (.doc)"""
        try:
            # 使用python-docx2txt处理经典Word文档
            # 注意：对于真正的.doc格式，需要使用antiword或python-docx2txt的扩展功能
            try:
                import subprocess
                import tempfile
                import os

                # 尝试使用antiword工具（如果可用）
                antiword_cmd = ['antiword', str(path)]
                result = subprocess.run(antiword_cmd, capture_output=True, text=True, timeout=30)

                if result.returncode == 0:
                    text = result.stdout.strip()
                    logger.info(f"经典Word文档解析成功: {path.name}, 文本长度: {len(text)}字符")

                    return ParsedContent(
                        text=text,
                        title=path.stem,
                        language="zh" if self._is_chinese_text(text) else "en",
                        confidence=0.8,
                        metadata={
                            "format": "doc",
                            "file_extension": ".doc",
                            "file_size": path.stat().st_size,
                            "parser": "antiword"
                        }
                    )
                else:
                    logger.warning(f"antiword解析失败: {result.stderr}")
                    raise Exception("antiword工具不可用")

            except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
                logger.warning(f"antiword工具不可用或解析失败: {e}")

                # 降级处理：仅提取元数据
                return ParsedContent(
                    text=f"[经典Word文档] - 内容提取功能暂不可用，请安装antiword工具",
                    title=path.stem,
                    language="metadata",
                    confidence=0.3,
                    metadata={
                        "format": "doc",
                        "file_extension": ".doc",
                        "file_size": path.stat().st_size,
                        "parser": "fallback",
                        "note": "需要安装antiword工具以支持内容提取"
                    }
                )

        except Exception as e:
            logger.error(f"经典Word文档解析失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "doc", "error": str(e)}
            )

    def _parse_ppt(self, path: Path) -> ParsedContent:
        """解析经典PowerPoint演示文稿 (.ppt)"""
        try:
            # 对于经典.ppt格式，可以使用python-pptx的有限支持或建议转换为.pptx
            try:
                # 尝试使用python-pptx处理（有限支持）
                from pptx import Presentation

                presentation = Presentation(str(path))
                text_content = []

                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            # 清理PPT文本中的乱码字符（垂直制表符等）
                            cleaned_text = self._clean_pptx_text(shape.text)
                            if cleaned_text.strip():
                                text_content.append(cleaned_text.strip())

                text = '\n'.join(text_content)
                logger.info(f"经典PowerPoint解析成功: {path.name}, 文本长度: {len(text)}字符")

                return ParsedContent(
                    text=text,
                    title=path.stem,
                    language="zh" if self._is_chinese_text(text) else "en",
                    confidence=0.7,
                    metadata={
                        "format": "ppt",
                        "file_extension": ".ppt",
                        "file_size": path.stat().st_size,
                        "slides_count": len(presentation.slides),
                        "parser": "python-pptx"
                    }
                )

            except Exception as e:
                logger.warning(f"python-pptx解析经典PPT失败: {e}")

                # 降级处理
                return ParsedContent(
                    text=f"[经典PowerPoint演示文稿] - 内容提取功能有限，建议转换为.pptx格式",
                    title=path.stem,
                    language="metadata",
                    confidence=0.3,
                    metadata={
                        "format": "ppt",
                        "file_extension": ".ppt",
                        "file_size": path.stat().st_size,
                        "parser": "fallback",
                        "note": "建议转换为.pptx格式以获得更好的支持"
                    }
                )

        except Exception as e:
            logger.error(f"经典PowerPoint解析失败: {e}")
            return ParsedContent(
                text="",
                title=None,
                language="metadata",
                confidence=0.0,
                metadata={"format": "ppt", "error": str(e)}
            )

    def _is_chinese_text(self, text: str) -> bool:
        """检测文本是否主要为中文"""
        if not text:
            return False

        # 简化的中文检测逻辑
        chinese_chars = len([c for c in text if '\u4e00' <= c <= '\u9fff'])
        english_chars = len([c for c in text if c.isalpha() and c.isascii()])

        if chinese_chars > english_chars:
            return True
        elif english_chars > 0:
            return False
        else:
            return False

    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式"""
        return list(self.supported_formats.keys())

    async def _extract_text_from_image(self, path: Path) -> str:
        """
        使用PaddleOCR从图片中提取文字内容

        Args:
            path: 图片文件路径

        Returns:
            str: 提取的文字内容
        """
        try:
            return await self._ocr_with_paddle(path)
        except ImportError:
            logger.warning("PaddleOCR未安装，请运行: pip install paddlepaddle paddleocr")
            return ""
        except Exception as e:
            logger.error(f"PaddleOCR文字提取失败: {str(e)}")
            return ""

    async def _ocr_with_paddle(self, path: Path) -> str:
        """使用PaddleOCR进行文字识别"""
        try:
            from paddleocr import PaddleOCR
        except ImportError:
            raise ImportError("PaddleOCR未安装，请运行: pip install paddlepaddle paddleocr")

        global _paddle_ocr_instance, _paddle_ocr_lock

        # 确保OCR实例已初始化
        async with _paddle_ocr_lock:
            if _paddle_ocr_instance is None:
                logger.info("初始化PaddleOCR实例...")
                try:
                    # 初始化PaddleOCR，使用中英文识别模型
                    _paddle_ocr_instance = PaddleOCR(
                        use_angle_cls=True,  # 使用角度分类
                        lang='ch',  # 支持中英文
                        det_db_thresh=0.3,  # 检测阈值
                        rec_batch_num=6  # 批量识别数量
                    )
                    logger.info("PaddleOCR实例初始化成功")
                except Exception as e:
                    logger.error(f"PaddleOCR实例初始化失败: {str(e)}")
                    raise RuntimeError(f"PaddleOCR初始化失败: {str(e)}")

        # 验证文件存在
        if not path.exists():
            logger.warning(f"图片文件不存在: {path}")
            return ""

        # 在线程池中执行OCR
        loop = asyncio.get_event_loop()

        def paddle_ocr_sync():
            try:
                # 使用绝对路径执行OCR识别
                abs_path = str(path.resolve())
                logger.debug(f"开始OCR识别: {abs_path}")

                result = _paddle_ocr_instance.ocr(abs_path)

                # 详细记录OCR结果
                logger.debug(f"OCR原始结果类型: {type(result)}")
                logger.debug(f"OCR原始结果: {result}")

                texts = []

                # PaddleOCR返回格式可能是字典或列表
                if isinstance(result, list) and len(result) > 0 and isinstance(result[0], dict):
                    # 新格式：字典格式
                    ocr_data = result[0]
                    if 'rec_texts' in ocr_data and 'rec_scores' in ocr_data:
                        rec_texts = ocr_data['rec_texts']
                        rec_scores = ocr_data['rec_scores']

                        logger.debug(f"识别到 {len(rec_texts)} 个文本行")
                        for i, (text, score) in enumerate(zip(rec_texts, rec_scores)):
                            logger.debug(f"第{i+1}行: '{text}' (置信度: {score:.3f})")

                            if text.strip() and score > 0.3:  # 降低置信度阈值
                                texts.append(text.strip())
                                logger.debug(f"采用文字: '{text.strip()}' (置信度: {score:.2f})")
                    else:
                        logger.debug("新格式中未找到rec_texts或rec_scores")

                elif isinstance(result, list) and len(result) > 0 and result[0]:
                    # 旧格式：列表格式
                    logger.debug(f"检测到 {len(result[0])} 个文本行")
                    for i, line in enumerate(result[0]):
                        logger.debug(f"第{i+1}行: {line}")
                        if line and len(line) >= 2:
                            # line[1] 包含文字和置信度
                            text = line[1][0] if line[1] and len(line[1]) > 0 else ""
                            confidence = line[1][1] if line[1] and len(line[1]) > 1 else 0.0

                            logger.debug(f"文字: '{text}', 置信度: {confidence:.3f}")

                            # 过滤低置信度和空文字
                            if text.strip() and confidence > 0.3:  # 降低置信度阈值
                                texts.append(text.strip())
                                logger.debug(f"采用文字: '{text.strip()}' (置信度: {confidence:.2f})")
                else:
                    logger.debug("OCR未检测到任何文本区域或格式不匹配")

                recognized_text = " ".join(texts)
                logger.debug(f"最终OCR结果: '{recognized_text}' (总文字数: {len(recognized_text)})")
                return recognized_text

            except Exception as e:
                logger.error(f"OCR识别过程中出错: {str(e)}")
                return ""

        return await loop.run_in_executor(None, paddle_ocr_sync)

    

# 测试代码
if __name__ == "__main__":
    # 配置日志
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )

    # 测试内容解析
    parser = ContentParser()

    # 测试当前目录的文件
    for file_path in Path(".").glob("*"):
        if file_path.is_file() and file_path.suffix in parser.get_supported_formats():
            print(f"\n解析文件内容: {file_path}")
            parsed = parser.parse_content(str(file_path))
            print(f"标题: {parsed.title}")
            print(f"语言: {parsed.language}")
            print(f"编码: {parsed.encoding}")
            print(f"置信度: {parsed.confidence}")
            print(f"内容长度: {len(parsed.text)} 字符")
            print(f"内容预览: {parsed.text[:200]}...")